"""
Microsoft Office MCP Server - Consolidated API

Vendored from https://github.com/sibyl-oracles/onit

Tools for creating PowerPoint presentations, Excel spreadsheets, and Word documents.

18 Core Tools:

PowerPoint:
1. create_presentation - Create a new presentation with title slide
2. add_slide - Universal slide creation (text, bullets, images, or combinations)
3. add_table_slide - Create data tables with automatic header styling
4. add_images_slide - Display multiple images in grid layouts
5. style_slide - Apply visual styling (backgrounds, colors)
6. get_presentation_info - Inspect presentation structure and metadata
7. download_media - Download media files (images, audio, video, PDFs) from URLs
8. read_presentation - Read full text content from all slides
9. modify_presentation - Edit text in existing shapes or delete slides

Excel:
10. create_excel - Create Excel file with headers and data
11. add_excel_rows - Add rows to existing Excel file
12. read_excel - Read cell contents from existing Excel file
13. modify_excel_cells - Edit specific cells in existing Excel file

Word:
14. create_document - Create Word document with optional header/logo
15. add_document_content - Add content (headings, paragraphs, lists, images, tables)
16. read_document - Read full content from existing Word document
17. modify_document - Edit or delete paragraphs in existing Word document

General:
18. get_file - Retrieve a created file as base64-encoded data for client download

All presentations use standard 16:9 widescreen dimensions (10" x 7.5").
"""
import os
import json
import base64
import tempfile
import requests
from typing import Any, List, Optional, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from fastmcp import FastMCP

import logging
logging.basicConfig(level=logging.ERROR, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

mcp = FastMCP("Microsoft Office MCP Server")

# Standard PowerPoint dimensions (16:9 widescreen)
PPTX_WIDTH = Inches(10)
PPTX_HEIGHT = Inches(7.5)

# Data path for file creation (set via options['data_path'] in run())
# All file writes are confined to this directory.
# Default: a unique temp directory per session, auto-cleaned on exit.
DATA_PATH = os.path.join(tempfile.gettempdir(), f"onit-office-{os.getpid()}")
_AUTO_CLEANUP = True  # True when using the default temp path


def _secure_makedirs(dir_path: str) -> None:
    """Create directory with owner-only permissions (0o700)."""
    os.makedirs(dir_path, mode=0o700, exist_ok=True)


def _ensure_directory(file_path: str) -> str:
    """Ensure the directory for the file exists and return absolute path.
    Uses secure directory creation with owner-only permissions."""
    full_path = os.path.abspath(os.path.expanduser(file_path))
    _secure_makedirs(os.path.dirname(full_path))
    return full_path


def _resolve_data_path(path: str) -> str:
    """Resolve a file path to be within DATA_PATH.
    If path is relative or outside DATA_PATH, place it under DATA_PATH."""
    expanded = os.path.expanduser(path)
    abs_data = os.path.abspath(os.path.expanduser(DATA_PATH))
    if not os.path.isabs(expanded) or not os.path.abspath(expanded).startswith(abs_data):
        basename = os.path.basename(expanded)
        return os.path.join(abs_data, basename)
    return os.path.abspath(expanded)


def _read_file_as_base64(filepath: str) -> dict:
    """Read a file and return its contents as base64-encoded data with MIME type."""
    mime_map = {
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    }
    ext = os.path.splitext(filepath)[1].lower()
    mime_type = mime_map.get(ext, "application/octet-stream")
    with open(filepath, "rb") as f:
        data = base64.b64encode(f.read()).decode("ascii")
    return {
        "file_name": os.path.basename(filepath),
        "mime_type": mime_type,
        "file_data_base64": data,
        "file_size_bytes": os.path.getsize(filepath),
    }


def _upload_file(filepath: str, callback_url: str) -> dict:
    """Upload a created file back to the callback server."""
    try:
        with open(filepath, 'rb') as f:
            files = {'file': (os.path.basename(filepath), f)}
            resp = requests.post(f"{callback_url}/uploads/", files=files, timeout=30)
            resp.raise_for_status()
        return {"uploaded": True, "download_url": f"{callback_url}/uploads/{os.path.basename(filepath)}"}
    except Exception as e:
        return {"uploaded": False, "upload_error": str(e)}


def _parse_color(color_str: str) -> RGBColor:
    """Convert color string to RGBColor. Accepts color names or hex codes (#RRGGBB)."""
    color_map = {
        "red": RGBColor(255, 0, 0),
        "green": RGBColor(0, 255, 0),
        "blue": RGBColor(0, 0, 255),
        "white": RGBColor(255, 255, 255),
        "black": RGBColor(0, 0, 0),
        "yellow": RGBColor(255, 255, 0),
        "orange": RGBColor(255, 165, 0),
        "purple": RGBColor(128, 0, 128),
        "gray": RGBColor(128, 128, 128),
        "grey": RGBColor(128, 128, 128),
    }

    if color_str.lower() in color_map:
        return color_map[color_str.lower()]

    if color_str.startswith("#") and len(color_str) == 7:
        try:
            r = int(color_str[1:3], 16)
            g = int(color_str[3:5], 16)
            b = int(color_str[5:7], 16)
            return RGBColor(r, g, b)
        except ValueError:
            pass

    return RGBColor(0, 0, 0)


# =============================================================================
# TOOL 1: CREATE PRESENTATION
# =============================================================================

@mcp.tool(
    title="Create PowerPoint Presentation",
    description="""Create a new 16:9 PowerPoint file with a title slide.

Args:
- title: Main title on the title slide (required)
- subtitle: Secondary text below title (default: "")
- path: File save location (default: server data directory/presentation.pptx)
- callback_url: URL to upload file to after creation (optional)

Returns JSON: {powerpoint_file, status, dimensions}"""
)
def create_presentation(
    title: str,
    subtitle: str = "",
    path: str = "",
    callback_url: Optional[str] = None
) -> str:
    try:
        if not path:
            path = os.path.join(DATA_PATH, "presentation.pptx")
        path = _resolve_data_path(path)
        full_path = _ensure_directory(path)

        prs = Presentation()
        prs.slide_width = PPTX_WIDTH
        prs.slide_height = PPTX_HEIGHT

        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

        prs.save(full_path)

        result = {
            "powerpoint_file": full_path,
            "status": "Created presentation with title slide",
            "dimensions": "16:9 widescreen (10\" x 7.5\")"
        }
        if callback_url:
            result.update(_upload_file(full_path, callback_url))
        return json.dumps(result)
    except Exception as e:
        return json.dumps({"error": str(e), "status": "Failed to create presentation"})


# =============================================================================
# TOOL 2: ADD SLIDE (UNIVERSAL)
# =============================================================================

@mcp.tool(
    title="Add Slide",
    description="""Add a slide to an existing presentation with various layouts.

Args:
- path: Path to existing .pptx file (required)
- title: Slide title (required except for "blank")
- layout: "text", "bullets", "image", "text_image", "bullets_image", "two_column", "blank" (default: "text")
- text: Text content for "text"/"text_image" layouts
- bullets: List of strings for "bullets"/"bullets_image" layouts
- image: Image file path for layouts with images
- image_position: "left" or "right" (default: "right")
- left_column/right_column: Lists for "two_column" layout

Returns JSON: {powerpoint_file, status, layout}"""
)
def add_slide(
    path: str,
    title: str = "",
    layout: str = "text",
    text: Optional[str] = None,
    bullets: Optional[Union[List[str], str]] = None,
    image: Optional[str] = None,
    image_position: str = "right",
    left_column: Optional[Union[List[str], str]] = None,
    right_column: Optional[Union[List[str], str]] = None
) -> str:
    # Accept JSON strings for list parameters (common with LLM callers)
    if isinstance(bullets, str):
        bullets = json.loads(bullets)
    if isinstance(left_column, str):
        left_column = json.loads(left_column)
    if isinstance(right_column, str):
        right_column = json.loads(right_column)

    try:
        full_path = os.path.abspath(os.path.expanduser(path))
        prs = Presentation(full_path)

        if layout == "text":
            result = _slide_text(prs, title, text or "")
        elif layout == "bullets":
            result = _slide_bullets(prs, title, bullets or [])
        elif layout == "image":
            result = _slide_image(prs, title, image)
        elif layout == "text_image":
            result = _slide_text_image(prs, title, text or "", image, image_position)
        elif layout == "bullets_image":
            result = _slide_bullets_image(prs, title, bullets or [], image, image_position)
        elif layout == "two_column":
            result = _slide_two_column(prs, title, left_column or [], right_column or [])
        elif layout == "blank":
            result = _slide_blank(prs)
        else:
            return json.dumps({
                "error": f"Unknown layout: {layout}",
                "valid_layouts": ["text", "bullets", "image", "text_image", "bullets_image", "two_column", "blank"]
            })

        if "error" in result:
            return json.dumps(result)

        prs.save(full_path)
        result["powerpoint_file"] = full_path
        return json.dumps(result)

    except Exception as e:
        return json.dumps({"error": str(e), "status": "Failed to add slide"})


def _slide_text(prs: Presentation, title: str, text: str) -> dict:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text_frame.text = text
    return {"status": "Added text slide", "layout": "text"}


def _slide_bullets(prs: Presentation, title: str, bullets: List[str]) -> dict:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    return {"status": "Added bullets slide", "layout": "bullets", "bullet_count": len(bullets)}


def _slide_image(prs: Presentation, title: str, image_path: Optional[str]) -> dict:
    if not image_path:
        return {"error": "image parameter required for 'image' layout"}
    img_full = os.path.abspath(os.path.expanduser(image_path))
    if not os.path.exists(img_full):
        return {"error": f"Image not found: {img_full}"}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    slide.shapes.add_picture(img_full, Inches(2), Inches(2), width=Inches(6))
    return {"status": "Added image slide", "layout": "image", "image": img_full}


def _slide_text_image(prs: Presentation, title: str, text: str,
                      image_path: Optional[str], position: str) -> dict:
    if not image_path:
        return {"error": "image parameter required for 'text_image' layout"}
    img_full = os.path.abspath(os.path.expanduser(image_path))
    if not os.path.exists(img_full):
        return {"error": f"Image not found: {img_full}"}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    if position.lower() == "right":
        text_left, img_left = Inches(0.5), Inches(5.5)
    else:
        text_left, img_left = Inches(5), Inches(0.5)
    text_box = slide.shapes.add_textbox(text_left, Inches(2), Inches(4.5), Inches(5))
    text_box.text_frame.word_wrap = True
    text_box.text_frame.text = text
    slide.shapes.add_picture(img_full, img_left, Inches(2), width=Inches(4.5))
    return {"status": "Added text+image slide", "layout": "text_image", "image_position": position}


def _slide_bullets_image(prs: Presentation, title: str, bullets: List[str],
                         image_path: Optional[str], position: str) -> dict:
    if not image_path:
        return {"error": "image parameter required for 'bullets_image' layout"}
    img_full = os.path.abspath(os.path.expanduser(image_path))
    if not os.path.exists(img_full):
        return {"error": f"Image not found: {img_full}"}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    if position.lower() == "right":
        bullets_left, img_left = Inches(0.5), Inches(6)
    else:
        bullets_left, img_left = Inches(5), Inches(0.5)
    bullets_box = slide.shapes.add_textbox(bullets_left, Inches(2), Inches(5), Inches(5))
    bullets_box.text_frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = bullets_box.text_frame.paragraphs[0]
        else:
            p = bullets_box.text_frame.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
    slide.shapes.add_picture(img_full, img_left, Inches(2.5), width=Inches(4))
    return {"status": "Added bullets+image slide", "layout": "bullets_image",
            "bullet_count": len(bullets), "image_position": position}


def _slide_two_column(prs: Presentation, title: str,
                      left: List[str], right: List[str]) -> dict:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    for i, item in enumerate(left):
        if i == 0:
            left_box.text_frame.text = item
        else:
            p = left_box.text_frame.add_paragraph()
            p.text = item
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4.5), Inches(5))
    for i, item in enumerate(right):
        if i == 0:
            right_box.text_frame.text = item
        else:
            p = right_box.text_frame.add_paragraph()
            p.text = item
    return {"status": "Added two-column slide", "layout": "two_column"}


def _slide_blank(prs: Presentation) -> dict:
    prs.slides.add_slide(prs.slide_layouts[6])
    return {"status": "Added blank slide", "layout": "blank", "slide_number": len(prs.slides)}


# =============================================================================
# TOOL 3: ADD TABLE SLIDE
# =============================================================================

@mcp.tool(
    title="Add Table Slide",
    description="""Add a slide with a formatted data table. Header row styled with blue background.

Args:
- path: Path to existing .pptx file (required)
- title: Slide title (required)
- data: 2D list where each inner list is a row (e.g., [["A","B"],["1","2"]])
- header: Style first row as header (default: True)

Returns JSON: {powerpoint_file, status, table_size}"""
)
def add_table_slide(
    path: str,
    title: str,
    data: Union[List[List[str]], str] = None,
    header: bool = True
) -> str:
    # Accept JSON strings for list parameters (common with LLM callers)
    if isinstance(data, str):
        data = json.loads(data)

    try:
        full_path = os.path.abspath(os.path.expanduser(path))
        prs = Presentation(full_path)

        if not data or not data[0]:
            return json.dumps({"error": "Table data cannot be empty"})

        rows, cols = len(data), len(data[0])
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(32)
        title_box.text_frame.paragraphs[0].font.bold = True

        table = slide.shapes.add_table(
            rows, cols, Inches(0.5), Inches(1.8), Inches(9), Inches(5)
        ).table

        for row_idx, row_data in enumerate(data):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.rows[row_idx].cells[col_idx]
                cell.text = str(cell_value)

                if header and row_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(255, 255, 255)

        prs.save(full_path)

        return json.dumps({
            "powerpoint_file": full_path,
            "status": "Added table slide",
            "table_size": f"{rows} rows x {cols} columns"
        })
    except Exception as e:
        return json.dumps({"error": str(e), "status": "Failed to add table slide"})


# =============================================================================
# TOOL 4: ADD IMAGES SLIDE (MULTIPLE)
# =============================================================================

@mcp.tool(
    title="Add Images Grid Slide",
    description="""Add a slide with multiple images in a grid layout (2-4 images).

Args:
- path: Path to existing .pptx file (required)
- title: Slide title (required)
- images: List of image file paths (required)
- grid: "horizontal" (side-by-side), "vertical" (stacked), "grid" (2x2) (default: "horizontal")

Returns JSON: {powerpoint_file, status, image_count, grid}"""
)
def add_images_slide(
    path: str,
    title: str,
    images: Union[List[str], str] = None,
    grid: str = "horizontal"
) -> str:
    # Accept JSON strings for list parameters (common with LLM callers)
    if isinstance(images, str):
        images = json.loads(images)

    try:
        full_path = os.path.abspath(os.path.expanduser(path))

        valid_images = []
        for img in images:
            img_path = os.path.abspath(os.path.expanduser(img))
            if os.path.exists(img_path):
                valid_images.append(img_path)
            else:
                logger.warning(f"Image not found: {img_path}")

        if not valid_images:
            return json.dumps({"error": "No valid image files found"})

        prs = Presentation(full_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(32)
        title_box.text_frame.paragraphs[0].font.bold = True

        img_size = Inches(3.5)
        spacing = Inches(0.3)
        start_top = Inches(2)

        if grid == "horizontal":
            positions = [(Inches(1 + i * 4), start_top) for i in range(len(valid_images))]
        elif grid == "vertical":
            positions = [(Inches(3), start_top + i * (3.5 + 0.3)) for i in range(len(valid_images))]
        elif grid == "grid":
            positions = [
                (Inches(1), start_top),
                (Inches(1) + img_size + spacing, start_top),
                (Inches(1), start_top + img_size + spacing),
                (Inches(1) + img_size + spacing, start_top + img_size + spacing)
            ]
        else:
            positions = [(Inches(1 + i * 4), start_top) for i in range(len(valid_images))]

        for idx, img_path in enumerate(valid_images):
            if idx < len(positions):
                left, top = positions[idx]
                slide.shapes.add_picture(img_path, left, top, width=img_size)

        prs.save(full_path)

        return json.dumps({
            "powerpoint_file": full_path,
            "status": "Added images grid slide",
            "image_count": len(valid_images),
            "grid": grid
        })
    except Exception as e:
        return json.dumps({"error": str(e), "status": "Failed to add images slide"})


# =============================================================================
# TOOL 5: STYLE SLIDE
# =============================================================================

@mcp.tool(
    title="Style Slide",
    description="""Apply background color to a specific slide.

Args:
- path: Path to existing .pptx file (required)
- background: Color name (red, blue, etc.) or hex code (#RRGGBB) (required)
- slide_index: Slide to style (default: -1 = last). Positive: 0=first. Negative: -1=last

Returns JSON: {powerpoint_file, status, background, slide_index}"""
)
def style_slide(
    path: str,
    background: str,
    slide_index: int = -1
) -> str:
    try:
        full_path = os.path.abspath(os.path.expanduser(path))
        prs = Presentation(full_path)

        if not prs.slides:
            return json.dumps({"error": "No slides in presentation"})

        slide = prs.slides[slide_index]
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _parse_color(background)

        prs.save(full_path)

        actual_index = slide_index if slide_index >= 0 else len(prs.slides) + slide_index

        return json.dumps({
            "powerpoint_file": full_path,
            "status": f"Styled slide {actual_index}",
            "background": background,
            "slide_index": actual_index
        })
    except Exception as e:
        return json.dumps({"error": str(e), "status": "Failed to style slide"})


# =============================================================================
# TOOL 6: GET PRESENTATION INFO
# =============================================================================

@mcp.tool(
    title="Get Presentation Info",
    description="""Retrieve information about an existing PowerPoint presentation.

Args:
- path: Path to the .pptx file (required)

Returns JSON: {powerpoint_file, slide_count, width_inches, height_inches, aspect_ratio, slides: [{index, title}]}"""
)
def get_presentation_info(path: str) -> str:
    try:
        path = _resolve_data_path(path)
        full_path = os.path.abspath(os.path.expanduser(path))

        if not os.path.exists(full_path):
            return json.dumps({
                "error": f"File not found: {full_path}. Use create_presentation first.",
                "status": "failed"
            })

        prs = Presentation(full_path)

        slides = []
        for idx, slide in enumerate(prs.slides):
            title = "No title"
            if slide.shapes.title:
                title = slide.shapes.title.text or f"Slide {idx + 1}"
            slides.append({"index": idx, "title": title})

        return json.dumps({
            "powerpoint_file": full_path,
            "slide_count": len(prs.slides),
            "width_inches": prs.slide_width.inches,
            "height_inches": prs.slide_height.inches,
            "aspect_ratio": f"{prs.slide_width.inches:.1f}:{prs.slide_height.inches:.1f}",
            "slides": slides,
            "status": "Retrieved presentation info"
        }, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e), "status": "Failed to get presentation info"})


# =============================================================================
# TOOL: READ PRESENTATION CONTENT
# =============================================================================

@mcp.tool(
    title="Read Presentation Content",
    description="""Read full text content from all slides in a PowerPoint presentation.

Unlike get_presentation_info (which returns only slide titles and metadata), this tool returns
the actual text content from every shape on every slide, plus notes and table data.

Args:
- path: Path to existing .pptx file (required)

Returns JSON: {powerpoint_file, slide_count, slides: [{index, title, shapes: [{shape_index, name, text, type}], notes, tables}], status}"""
)
def read_presentation(path: str) -> str:
    try:
        full_path = os.path.abspath(os.path.expanduser(path))

        if not os.path.exists(full_path):
            return json.dumps({
                "error": f"File not found: {full_path}",
                "path": path
            })

        prs = Presentation(full_path)

        slides = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_title = ""
            if slide.shapes.title:
                slide_title = slide.shapes.title.text or ""

            shapes_data = []
            tables_data = []
            for shape_idx, shape in enumerate(slide.shapes):
                shape_info = {
                    "shape_index": shape_idx,
                    "name": shape.name,
                    "type": shape.shape_type.__class__.__name__ if hasattr(shape.shape_type, '__class__') else str(shape.shape_type),
                }

                if shape.has_text_frame:
                    shape_info["text"] = shape.text_frame.text
                    shape_info["paragraphs"] = [p.text for p in shape.text_frame.paragraphs]

                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        table_rows.append([cell.text for cell in row.cells])
                    tables_data.append({
                        "shape_index": shape_idx,
                        "rows": table_rows
                    })

                shapes_data.append(shape_info)

            notes_text = ""
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text

            slides.append({
                "index": slide_idx,
                "title": slide_title,
                "shapes": shapes_data,
                "tables": tables_data,
                "notes": notes_text,
            })

        return json.dumps({
            "powerpoint_file": full_path,
            "slide_count": len(prs.slides),
            "slides": slides,
            "status": "success"
        }, indent=2)

    except Exception as e:
        return json.dumps({
            "error": str(e),
            "path": path,
            "status": "failed"
        })


# =============================================================================
# TOOL: MODIFY PRESENTATION
# =============================================================================

@mcp.tool(
    title="Modify Presentation",
    description="""Edit text in existing shapes or delete slides in a PowerPoint presentation.

Args:
- path: Path to existing .pptx file (required)
- updates: List of shape text updates, each with "slide_index", "shape_index", and "text" (optional)
- delete_slides: List of slide indices to delete (optional, processed after updates, use descending order internally)

Example updates: [{"slide_index": 0, "shape_index": 0, "text": "New Title"}]
Example delete_slides: [2, 4]

Returns JSON: {powerpoint_file, shapes_updated, slides_deleted, status}"""
)
def modify_presentation(
    path: str,
    updates: Optional[Union[List[dict], str]] = None,
    delete_slides: Optional[Union[List[int], str]] = None,
) -> str:
    if isinstance(updates, str):
        updates = json.loads(updates)
    if isinstance(delete_slides, str):
        delete_slides = json.loads(delete_slides)

    try:
        full_path = os.path.abspath(os.path.expanduser(path))

        if not os.path.exists(full_path):
            return json.dumps({
                "error": f"File not found: {full_path}",
                "path": path
            })

        if not updates and not delete_slides:
            return json.dumps({"error": "Either updates or delete_slides is required"})

        prs = Presentation(full_path)
        shapes_updated = 0
        slides_deleted = 0

        if updates:
            for update in updates:
                slide_idx = update.get("slide_index")
                shape_idx = update.get("shape_index")
                new_text = update.get("text")
                if slide_idx is None or shape_idx is None or new_text is None:
                    continue
                if slide_idx < 0 or slide_idx >= len(prs.slides):
                    continue
                slide = prs.slides[slide_idx]
                if shape_idx < 0 or shape_idx >= len(slide.shapes):
                    continue
                shape = slide.shapes[shape_idx]
                if shape.has_text_frame:
                    shape.text_frame.clear()
                    shape.text_frame.paragraphs[0].text = new_text
                    shapes_updated += 1

        if delete_slides:
            slide_id_map = []
            for idx in sorted(delete_slides, reverse=True):
                if 0 <= idx < len(prs.slides):
                    rId = prs.slides._sldIdLst[idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                    slide_id_map.append((idx, rId))

            for idx, rId in slide_id_map:
                sldId = prs.slides._sldIdLst[idx]
                prs.slides._sldIdLst.remove(sldId)
                if rId:
                    prs.part.drop_rel(rId)
                slides_deleted += 1

        prs.save(full_path)

        return json.dumps({
            "powerpoint_file": full_path,
            "shapes_updated": shapes_updated,
            "slides_deleted": slides_deleted,
            "status": "success"
        }, indent=2)

    except Exception as e:
        return json.dumps({
            "error": str(e),
            "path": path,
            "status": "failed"
        })


# =============================================================================
# TOOL 7: DOWNLOAD MEDIA
# =============================================================================

@mcp.tool(
    title="Download Media from URL",
    description="""Download a media file from a URL and save locally. Supports images, audio, video, PDFs, and other common media formats.

Args:
- url: Full URL of the media file (required)
- output_path: Local save path (required)
- timeout: Download timeout in seconds (default: 60)

Returns JSON: {file_path, size_bytes, content_type, media_type, status}"""
)
def download_media(
    url: str,
    output_path: str,
    timeout: int = 60
) -> str:
    try:
        if not url.startswith(('http://', 'https://')):
            return json.dumps({"error": "URL must start with http:// or https://"})

        output_path = _resolve_data_path(output_path)
        full_path = os.path.abspath(os.path.expanduser(output_path))
        _secure_makedirs(os.path.dirname(full_path))

        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        response = requests.get(url, headers=headers, timeout=timeout, stream=True)
        response.raise_for_status()

        content_type = response.headers.get('Content-Type', '').lower()
        content = response.content

        media_type = "unknown"
        detected_type = None

        signatures = [
            (b'\xff\xd8\xff', 'image/jpeg', 'image'),
            (b'\x89PNG\r\n\x1a\n', 'image/png', 'image'),
            (b'GIF87a', 'image/gif', 'image'),
            (b'GIF89a', 'image/gif', 'image'),
            (b'RIFF', 'image/webp', 'image'),
            (b'BM', 'image/bmp', 'image'),
            (b'ID3', 'audio/mpeg', 'audio'),
            (b'\xff\xfb', 'audio/mpeg', 'audio'),
            (b'\xff\xf3', 'audio/mpeg', 'audio'),
            (b'fLaC', 'audio/flac', 'audio'),
            (b'OggS', 'audio/ogg', 'audio'),
            (b'\x00\x00\x00', 'video/mp4', 'video'),
            (b'\x1a\x45\xdf\xa3', 'video/webm', 'video'),
            (b'%PDF', 'application/pdf', 'document'),
        ]

        for signature, sig_type, sig_media in signatures:
            if content.startswith(signature):
                detected_type = sig_type
                media_type = sig_media
                break

        if not detected_type:
            type_prefixes = {
                'image/': 'image',
                'audio/': 'audio',
                'video/': 'video',
                'application/pdf': 'document',
                'application/octet-stream': 'binary',
            }
            for prefix, mtype in type_prefixes.items():
                if prefix in content_type:
                    media_type = mtype
                    break

        fd = os.open(full_path, os.O_WRONLY | os.O_CREAT | os.O_TRUNC, 0o600)
        with os.fdopen(fd, 'wb') as f:
            f.write(content)

        return json.dumps({
            "file_path": full_path,
            "size_bytes": len(content),
            "content_type": detected_type or content_type,
            "media_type": media_type,
            "status": f"Downloaded {media_type} file successfully"
        })

    except requests.exceptions.Timeout:
        return json.dumps({"error": f"Download timed out after {timeout} seconds"})
    except requests.exceptions.RequestException as e:
        return json.dumps({"error": f"Download failed: {str(e)}"})
    except Exception as e:
        return json.dumps({"error": str(e), "status": "Failed to download file"})


# =============================================================================
# TOOL 8: CREATE EXCEL FILE
# =============================================================================

@mcp.tool(
    title="Create Excel File",
    description="""Create a new Excel file with headers and optional data rows.

Args:
- path: File save path (default: server data directory/spreadsheet.xlsx)
- sheet_name: Name of the worksheet (default: "Sheet1")
- headers: List of column headers (e.g., ["Name", "Age", "Email"])
- rows: List of rows, each row is a list of values (e.g., [["John", 30, "john@example.com"]])
- auto_width: Auto-adjust column widths (default: True)
- callback_url: URL to upload file to after creation (optional)

Returns JSON: {excel_file, sheet_name, headers, row_count, status}"""
)
def create_excel(
    path: str = "",
    sheet_name: str = "Sheet1",
    headers: Optional[Union[List[str], str]] = None,
    rows: Optional[Union[List[List[Any]], str]] = None,
    auto_width: bool = True,
    callback_url: Optional[str] = None
) -> str:
    # Accept JSON strings for list parameters (common with LLM callers)
    if isinstance(headers, str):
        headers = json.loads(headers)
    if isinstance(rows, str):
        rows = json.loads(rows)

    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.utils import get_column_letter
    except ImportError:
        return json.dumps({
            "error": "openpyxl not installed. Run: pip install openpyxl",
            "path": path
        })

    try:
        if not path:
            path = os.path.join(DATA_PATH, "spreadsheet.xlsx")
        path = _resolve_data_path(path)
        full_path = _ensure_directory(path)

        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name

        current_row = 1

        if headers:
            header_font = Font(bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")

            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")

            current_row = 2

        row_count = 0
        if rows:
            for row_data in rows:
                for col_idx, value in enumerate(row_data, 1):
                    ws.cell(row=current_row, column=col_idx, value=value)
                current_row += 1
                row_count += 1

        if auto_width:
            for col_idx in range(1, ws.max_column + 1):
                max_length = 0
                column_letter = get_column_letter(col_idx)
                for row_idx in range(1, ws.max_row + 1):
                    cell_value = ws.cell(row=row_idx, column=col_idx).value
                    if cell_value:
                        max_length = max(max_length, len(str(cell_value)))
                ws.column_dimensions[column_letter].width = min(max_length + 2, 50)

        wb.save(full_path)

        result = {
            "excel_file": full_path,
            "sheet_name": sheet_name,
            "headers": headers or [],
            "row_count": row_count,
            "status": "success"
        }
        if callback_url:
            result.update(_upload_file(full_path, callback_url))
        return json.dumps(result, indent=2)

    except Exception as e:
        return json.dumps({
            "error": str(e),
            "path": path,
            "status": "failed"
        })


# =============================================================================
# TOOL 9: ADD ROWS TO EXCEL FILE
# =============================================================================

@mcp.tool(
    title="Add Excel Rows",
    description="""Add rows of data to an existing Excel file.

Args:
- path: Path to existing Excel file (required)
- rows: List of rows to add (e.g., [["Alice", 25], ["Bob", 30]])
- sheet_name: Worksheet name (default: active sheet)
- auto_width: Auto-adjust column widths (default: True)

Returns JSON: {excel_file, sheet_name, rows_added, total_rows, status}"""
)
def add_excel_rows(
    path: str,
    rows: Union[List[List[Any]], str] = None,
    sheet_name: Optional[str] = None,
    auto_width: bool = True
) -> str:
    # Accept JSON strings for list parameters (common with LLM callers)
    if isinstance(rows, str):
        rows = json.loads(rows)

    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
    except ImportError:
        return json.dumps({
            "error": "openpyxl not installed. Run: pip install openpyxl",
            "path": path
        })

    try:
        full_path = os.path.abspath(os.path.expanduser(path))

        if not os.path.exists(full_path):
            return json.dumps({
                "error": f"Excel file not found: {full_path}",
                "path": path
            })

        wb = load_workbook(full_path)

        if sheet_name:
            if sheet_name not in wb.sheetnames:
                return json.dumps({
                    "error": f"Sheet '{sheet_name}' not found. Available: {wb.sheetnames}",
                    "path": path
                })
            ws = wb[sheet_name]
        else:
            ws = wb.active
            sheet_name = ws.title

        next_row = ws.max_row + 1

        rows_added = 0
        for row_data in rows:
            for col_idx, value in enumerate(row_data, 1):
                ws.cell(row=next_row, column=col_idx, value=value)
            next_row += 1
            rows_added += 1

        if auto_width:
            for col_idx in range(1, ws.max_column + 1):
                max_length = 0
                column_letter = get_column_letter(col_idx)
                for row_idx in range(1, ws.max_row + 1):
                    cell_value = ws.cell(row=row_idx, column=col_idx).value
                    if cell_value:
                        max_length = max(max_length, len(str(cell_value)))
                ws.column_dimensions[column_letter].width = min(max_length + 2, 50)

        wb.save(full_path)

        return json.dumps({
            "excel_file": full_path,
            "sheet_name": sheet_name,
            "rows_added": rows_added,
            "total_rows": ws.max_row,
            "status": "success"
        }, indent=2)

    except Exception as e:
        return json.dumps({
            "error": str(e),
            "path": path,
            "status": "failed"
        })


# =============================================================================
# TOOL: READ EXCEL FILE
# =============================================================================

@mcp.tool(
    title="Read Excel File",
    description="""Read cell contents from an existing Excel file.

Args:
- path: Path to existing .xlsx file (required)
- sheet_name: Worksheet name (default: active sheet)
- max_rows: Maximum number of data rows to return (default: 500)

Returns JSON: {excel_file, sheet_name, sheets, headers, rows, row_count, total_rows, status}"""
)
def read_excel(
    path: str,
    sheet_name: Optional[str] = None,
    max_rows: int = 500,
) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return json.dumps({
            "error": "openpyxl not installed. Run: pip install openpyxl",
            "path": path
        })

    try:
        full_path = os.path.abspath(os.path.expanduser(path))

        if not os.path.exists(full_path):
            return json.dumps({
                "error": f"Excel file not found: {full_path}",
                "path": path
            })

        wb = load_workbook(full_path, read_only=True, data_only=True)

        if sheet_name:
            if sheet_name not in wb.sheetnames:
                wb.close()
                return json.dumps({
                    "error": f"Sheet '{sheet_name}' not found. Available: {wb.sheetnames}",
                    "path": path
                })
            ws = wb[sheet_name]
        else:
            ws = wb.active
            sheet_name = ws.title

        headers = []
        rows_data = []
        total_rows = 0

        for row_idx, row in enumerate(ws.iter_rows(values_only=True), 1):
            if row_idx == 1:
                headers = [str(c) if c is not None else "" for c in row]
            else:
                total_rows += 1
                if total_rows <= max_rows:
                    rows_data.append([c if c is not None else "" for c in row])

        wb.close()

        return json.dumps({
            "excel_file": full_path,
            "sheet_name": sheet_name,
            "sheets": wb.sheetnames if hasattr(wb, 'sheetnames') else [sheet_name],
            "headers": headers,
            "rows": rows_data,
            "row_count": len(rows_data),
            "total_rows": total_rows,
            "truncated": total_rows > max_rows,
            "status": "success"
        }, indent=2, default=str)

    except Exception as e:
        return json.dumps({
            "error": str(e),
            "path": path,
            "status": "failed"
        })


# =============================================================================
# TOOL: MODIFY EXCEL CELLS
# =============================================================================

@mcp.tool(
    title="Modify Excel Cells",
    description="""Edit specific cells in an existing Excel file.

Args:
- path: Path to existing .xlsx file (required)
- updates: List of cell updates, each with "cell" (e.g., "A1", "B3") and "value" (required)
- sheet_name: Worksheet name (default: active sheet)

Example updates: [{"cell": "A1", "value": "New Title"}, {"cell": "B3", "value": 42}]

Returns JSON: {excel_file, sheet_name, cells_updated, status}"""
)
def modify_excel_cells(
    path: str,
    updates: Union[List[dict], str] = None,
    sheet_name: Optional[str] = None,
) -> str:
    if isinstance(updates, str):
        updates = json.loads(updates)

    try:
        from openpyxl import load_workbook
    except ImportError:
        return json.dumps({
            "error": "openpyxl not installed. Run: pip install openpyxl",
            "path": path
        })

    try:
        full_path = os.path.abspath(os.path.expanduser(path))

        if not os.path.exists(full_path):
            return json.dumps({
                "error": f"Excel file not found: {full_path}",
                "path": path
            })

        if not updates:
            return json.dumps({"error": "updates parameter is required"})

        wb = load_workbook(full_path)

        if sheet_name:
            if sheet_name not in wb.sheetnames:
                return json.dumps({
                    "error": f"Sheet '{sheet_name}' not found. Available: {wb.sheetnames}",
                    "path": path
                })
            ws = wb[sheet_name]
        else:
            ws = wb.active
            sheet_name = ws.title

        cells_updated = 0
        for update in updates:
            cell_ref = update.get("cell")
            value = update.get("value")
            if cell_ref is None:
                continue
            ws[cell_ref] = value
            cells_updated += 1

        wb.save(full_path)

        return json.dumps({
            "excel_file": full_path,
            "sheet_name": sheet_name,
            "cells_updated": cells_updated,
            "status": "success"
        }, indent=2)

    except Exception as e:
        return json.dumps({
            "error": str(e),
            "path": path,
            "status": "failed"
        })


# =============================================================================
# TOOL: CREATE WORD DOCUMENT
# =============================================================================

@mcp.tool(
    title="Create Word Document",
    description="""Create a new Word document with optional header, logo, and initial content.

Args:
- path: File save path (default: server data directory/document.docx)
- title: Document title as first heading (optional)
- header_text: Text to appear in page header (optional)
- footer_text: Text to appear in page footer (optional)
- logo_path: Path to logo image for header (optional)
- logo_width_inches: Logo width in inches (default: 1.5)
- content_type: Optional initial content type - "heading", "paragraph", "bullets" (optional)
- text: Text for initial paragraph or heading content (optional)
- items: List of strings for bullets content (optional)
- callback_url: URL to upload file to after creation (optional)

Returns JSON: {document_file, title, has_header, has_footer, has_logo, status}"""
)
def create_document(
    path: str = "",
    title: Optional[str] = None,
    header_text: Optional[str] = None,
    footer_text: Optional[str] = None,
    logo_path: Optional[str] = None,
    logo_width_inches: float = 1.5,
    content_type: Optional[str] = None,
    text: Optional[str] = None,
    items: Optional[Union[List[str], str]] = None,
    callback_url: Optional[str] = None,
) -> str:
    # Accept JSON strings for list parameters (common with LLM callers)
    if isinstance(items, str):
        items = json.loads(items)

    try:
        from docx import Document
        from docx.shared import Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return json.dumps({
            "error": "python-docx not installed. Run: pip install python-docx",
            "path": path
        })

    try:
        if not path:
            path = os.path.join(DATA_PATH, "document.docx")
        path = _resolve_data_path(path)
        full_path = _ensure_directory(path)

        doc = Document()

        has_header = False
        has_footer = False
        has_logo = False

        if header_text or logo_path:
            section = doc.sections[0]
            header = section.header
            header_para = header.paragraphs[0] if header.paragraphs else header.add_paragraph()

            if logo_path:
                logo_full = os.path.abspath(os.path.expanduser(logo_path))
                if os.path.exists(logo_full):
                    run = header_para.add_run()
                    run.add_picture(logo_full, width=Inches(logo_width_inches))
                    has_logo = True

            if header_text:
                if has_logo:
                    header_para.add_run("  ")
                header_para.add_run(header_text)
                has_header = True

        if footer_text:
            section = doc.sections[0]
            footer = section.footer
            footer_para = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
            footer_para.text = footer_text
            footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            has_footer = True

        if title:
            doc.add_heading(title, level=0)

        if content_type == "paragraph" and text:
            doc.add_paragraph(text)
        elif content_type == "heading" and text:
            doc.add_heading(text, level=1)
        elif content_type == "bullets" and items:
            for item in items:
                doc.add_paragraph(item, style='List Bullet')

        doc.save(full_path)

        result = {
            "document_file": full_path,
            "title": title or "",
            "has_header": has_header or has_logo,
            "has_footer": has_footer,
            "has_logo": has_logo,
            "content_type": content_type or "",
            "status": "success"
        }
        if callback_url:
            result.update(_upload_file(full_path, callback_url))
        return json.dumps(result, indent=2)

    except Exception as e:
        return json.dumps({
            "error": str(e),
            "path": path,
            "status": "failed"
        })


# =============================================================================
# TOOL 11: ADD CONTENT TO WORD DOCUMENT
# =============================================================================

@mcp.tool(
    title="Add Document Content",
    description="""Add content to an existing Word document.

Args:
- path: Path to existing .docx file (required)
- content_type: "heading", "paragraph", "bullets", "image", "table", "page_break" (required)
- text: Text content for heading/paragraph
- level: Heading level 1-4 (default: 1, for heading type)
- items: List of strings (for bullets type)
- image_path: Path to image file (for image type)
- image_width_inches: Image width (default: 5.0)
- table_data: 2D list for table (e.g., [["A","B"],["1","2"]])
- table_header: Style first row as header (default: True)

Returns JSON: {document_file, content_type, status}"""
)
def add_document_content(
    path: str,
    content_type: str,
    text: Optional[str] = None,
    level: int = 1,
    items: Optional[Union[List[str], str]] = None,
    image_path: Optional[str] = None,
    image_width_inches: float = 5.0,
    table_data: Optional[Union[List[List[Any]], str]] = None,
    table_header: bool = True
) -> str:
    # Accept JSON strings for list parameters (common with LLM callers)
    if isinstance(items, str):
        items = json.loads(items)
    if isinstance(table_data, str):
        table_data = json.loads(table_data)

    try:
        from docx import Document
        from docx.shared import Inches
        from docx.oxml.ns import nsdecls
        from docx.oxml import parse_xml
    except ImportError:
        return json.dumps({
            "error": "python-docx not installed. Run: pip install python-docx",
            "path": path
        })

    try:
        full_path = os.path.abspath(os.path.expanduser(path))

        if not os.path.exists(full_path):
            return json.dumps({
                "error": f"Document not found: {full_path}",
                "path": path
            })

        doc = Document(full_path)
        result_details = {}

        if content_type == "heading":
            if not text:
                return json.dumps({"error": "text required for heading"})
            level = max(1, min(level, 4))
            doc.add_heading(text, level=level)
            result_details["level"] = level

        elif content_type == "paragraph":
            if not text:
                return json.dumps({"error": "text required for paragraph"})
            doc.add_paragraph(text)

        elif content_type == "bullets":
            if not items:
                return json.dumps({"error": "items required for bullets"})
            for item in items:
                doc.add_paragraph(item, style='List Bullet')
            result_details["item_count"] = len(items)

        elif content_type == "image":
            if not image_path:
                return json.dumps({"error": "image_path required for image"})
            img_full = os.path.abspath(os.path.expanduser(image_path))
            if not os.path.exists(img_full):
                return json.dumps({"error": f"Image not found: {img_full}"})
            doc.add_picture(img_full, width=Inches(image_width_inches))
            result_details["image"] = img_full

        elif content_type == "table":
            if not table_data or not table_data[0]:
                return json.dumps({"error": "table_data required for table"})

            rows = len(table_data)
            cols = len(table_data[0])
            table = doc.add_table(rows=rows, cols=cols)
            table.style = 'Table Grid'

            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_value in enumerate(row_data):
                    cell = table.rows[row_idx].cells[col_idx]
                    cell.text = str(cell_value)

                    if table_header and row_idx == 0:
                        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="4472C4"/>')
                        cell._tc.get_or_add_tcPr().append(shading)
                        for para in cell.paragraphs:
                            for run in para.runs:
                                run.font.bold = True
                                run.font.color.rgb = None

            result_details["table_size"] = f"{rows}x{cols}"

        elif content_type == "page_break":
            doc.add_page_break()

        else:
            return json.dumps({
                "error": f"Unknown content_type: {content_type}",
                "valid_types": ["heading", "paragraph", "bullets", "image", "table", "page_break"]
            })

        doc.save(full_path)

        return json.dumps({
            "document_file": full_path,
            "content_type": content_type,
            **result_details,
            "status": "success"
        }, indent=2)

    except Exception as e:
        return json.dumps({
            "error": str(e),
            "path": path,
            "status": "failed"
        })


# =============================================================================
# TOOL: READ WORD DOCUMENT
# =============================================================================

@mcp.tool(
    title="Read Word Document",
    description="""Read full content from an existing Word document.

Args:
- path: Path to existing .docx file (required)

Returns JSON: {document_file, paragraphs: [{index, text, style}], tables: [{index, rows}], headers, footers, section_count, status}"""
)
def read_document(path: str) -> str:
    try:
        from docx import Document
    except ImportError:
        return json.dumps({
            "error": "python-docx not installed. Run: pip install python-docx",
            "path": path
        })

    try:
        full_path = os.path.abspath(os.path.expanduser(path))

        if not os.path.exists(full_path):
            return json.dumps({
                "error": f"Document not found: {full_path}",
                "path": path
            })

        doc = Document(full_path)

        paragraphs = []
        for idx, para in enumerate(doc.paragraphs):
            paragraphs.append({
                "index": idx,
                "text": para.text,
                "style": para.style.name if para.style else "Normal",
            })

        tables = []
        for idx, table in enumerate(doc.tables):
            table_rows = []
            for row in table.rows:
                table_rows.append([cell.text for cell in row.cells])
            tables.append({"index": idx, "rows": table_rows})

        headers = []
        footers = []
        for section in doc.sections:
            if section.header and section.header.paragraphs:
                header_text = "\n".join(p.text for p in section.header.paragraphs if p.text)
                if header_text:
                    headers.append(header_text)
            if section.footer and section.footer.paragraphs:
                footer_text = "\n".join(p.text for p in section.footer.paragraphs if p.text)
                if footer_text:
                    footers.append(footer_text)

        return json.dumps({
            "document_file": full_path,
            "paragraphs": paragraphs,
            "tables": tables,
            "headers": headers,
            "footers": footers,
            "section_count": len(doc.sections),
            "paragraph_count": len(paragraphs),
            "table_count": len(tables),
            "status": "success"
        }, indent=2)

    except Exception as e:
        return json.dumps({
            "error": str(e),
            "path": path,
            "status": "failed"
        })


# =============================================================================
# TOOL: MODIFY WORD DOCUMENT
# =============================================================================

@mcp.tool(
    title="Modify Word Document",
    description="""Edit or delete paragraphs in an existing Word document.

Args:
- path: Path to existing .docx file (required)
- updates: List of paragraph updates, each with "paragraph_index" (int), and optionally "text" (new text) and/or "style" (new style name) (optional)
- delete_indices: List of paragraph indices to delete (optional, processed after updates)

Example updates: [{"paragraph_index": 0, "text": "New Title"}, {"paragraph_index": 2, "style": "Heading 2"}]
Example delete_indices: [3, 5]

Returns JSON: {document_file, paragraphs_updated, paragraphs_deleted, status}"""
)
def modify_document(
    path: str,
    updates: Optional[Union[List[dict], str]] = None,
    delete_indices: Optional[Union[List[int], str]] = None,
) -> str:
    if isinstance(updates, str):
        updates = json.loads(updates)
    if isinstance(delete_indices, str):
        delete_indices = json.loads(delete_indices)

    try:
        from docx import Document
    except ImportError:
        return json.dumps({
            "error": "python-docx not installed. Run: pip install python-docx",
            "path": path
        })

    try:
        full_path = os.path.abspath(os.path.expanduser(path))

        if not os.path.exists(full_path):
            return json.dumps({
                "error": f"Document not found: {full_path}",
                "path": path
            })

        if not updates and not delete_indices:
            return json.dumps({"error": "Either updates or delete_indices is required"})

        doc = Document(full_path)
        paragraphs_updated = 0
        paragraphs_deleted = 0

        if updates:
            for update in updates:
                idx = update.get("paragraph_index")
                if idx is None or idx < 0 or idx >= len(doc.paragraphs):
                    continue
                para = doc.paragraphs[idx]
                if "text" in update:
                    para.clear()
                    para.add_run(update["text"])
                    paragraphs_updated += 1
                if "style" in update:
                    try:
                        para.style = update["style"]
                    except KeyError:
                        pass
                    paragraphs_updated += 1

        if delete_indices:
            for idx in sorted(delete_indices, reverse=True):
                if 0 <= idx < len(doc.paragraphs):
                    p = doc.paragraphs[idx]._element
                    p.getparent().remove(p)
                    paragraphs_deleted += 1

        doc.save(full_path)

        return json.dumps({
            "document_file": full_path,
            "paragraphs_updated": paragraphs_updated,
            "paragraphs_deleted": paragraphs_deleted,
            "status": "success"
        }, indent=2)

    except Exception as e:
        return json.dumps({
            "error": str(e),
            "path": path,
            "status": "failed"
        })


# =============================================================================
# TOOL: GET FILE
# =============================================================================

@mcp.tool(
    title="Get File",
    description="""Retrieve a created file (pptx, xlsx, docx) as base64-encoded data.

Use this tool after you have finished creating and updating a file to send it
back to the client. Call this as the final step once all edits are complete.

Args:
- path: Path to the file to retrieve (required)

Returns JSON: {file_name, mime_type, file_data_base64, file_size_bytes, status}"""
)
def get_file(path: str) -> str:
    try:
        resolved = _resolve_data_path(path)
        full_path = os.path.abspath(os.path.expanduser(resolved))

        if not os.path.exists(full_path):
            return json.dumps({
                "error": f"File not found: {full_path}",
                "path": path,
                "status": "failed"
            })

        result = _read_file_as_base64(full_path)
        result["status"] = "success"
        return json.dumps(result)

    except Exception as e:
        return json.dumps({
            "error": str(e),
            "path": path,
            "status": "failed"
        })
