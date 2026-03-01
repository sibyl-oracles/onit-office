"""Tests for PowerPoint tools in mcp_server.py."""

import json
import os

import pytest
from pptx import Presentation

from onit_office.mcp_server import (
    create_presentation,
    add_slide,
    add_table_slide,
    add_images_slide,
    style_slide,
    get_presentation_info,
    read_presentation,
    modify_presentation,
)


class TestCreatePresentation:
    """Tests for create_presentation()."""

    def test_basic_creation(self, tmp_data_dir):
        result = json.loads(create_presentation(
            title="My Presentation",
            path=str(tmp_data_dir / "basic.pptx"),
        ))
        assert result["status"] == "Created presentation with title slide"
        assert os.path.exists(result["powerpoint_file"])
        assert "16:9" in result["dimensions"]

    def test_with_subtitle(self, tmp_data_dir):
        path = str(tmp_data_dir / "subtitle.pptx")
        result = json.loads(create_presentation(
            title="Title", subtitle="Subtitle", path=path,
        ))
        assert "error" not in result
        prs = Presentation(path)
        slide = prs.slides[0]
        assert slide.placeholders[1].text == "Subtitle"

    def test_default_path(self, tmp_data_dir):
        result = json.loads(create_presentation(title="Default Path"))
        assert "error" not in result
        assert os.path.exists(result["powerpoint_file"])

    def test_presentation_dimensions(self, tmp_data_dir):
        path = str(tmp_data_dir / "dims.pptx")
        create_presentation(title="Dims", path=path)
        prs = Presentation(path)
        assert prs.slide_width.inches == 10.0
        assert prs.slide_height.inches == 7.5

    def test_creates_single_slide(self, tmp_data_dir):
        path = str(tmp_data_dir / "single.pptx")
        create_presentation(title="Single", path=path)
        prs = Presentation(path)
        assert len(prs.slides) == 1


class TestAddSlide:
    """Tests for add_slide()."""

    def test_text_layout(self, pptx_file):
        result = json.loads(add_slide(
            path=pptx_file, title="Text Slide", layout="text", text="Hello World",
        ))
        assert result["status"] == "Added text slide"
        assert result["layout"] == "text"
        prs = Presentation(pptx_file)
        assert len(prs.slides) == 2

    def test_bullets_layout(self, pptx_file):
        result = json.loads(add_slide(
            path=pptx_file, title="Bullets", layout="bullets",
            bullets=["Item 1", "Item 2", "Item 3"],
        ))
        assert result["layout"] == "bullets"
        assert result["bullet_count"] == 3

    def test_bullets_json_string(self, pptx_file):
        result = json.loads(add_slide(
            path=pptx_file, title="Bullets JSON", layout="bullets",
            bullets='["A", "B", "C"]',
        ))
        assert result["bullet_count"] == 3

    def test_image_layout(self, pptx_file, sample_image):
        result = json.loads(add_slide(
            path=pptx_file, title="Image Slide", layout="image",
            image=sample_image,
        ))
        assert result["layout"] == "image"

    def test_image_layout_missing_image(self, pptx_file):
        result = json.loads(add_slide(
            path=pptx_file, title="Missing", layout="image",
        ))
        assert "error" in result

    def test_image_layout_nonexistent_file(self, pptx_file):
        result = json.loads(add_slide(
            path=pptx_file, title="Bad", layout="image",
            image="/nonexistent/image.png",
        ))
        assert "error" in result

    def test_text_image_layout(self, pptx_file, sample_image):
        result = json.loads(add_slide(
            path=pptx_file, title="Text+Image", layout="text_image",
            text="Description", image=sample_image, image_position="right",
        ))
        assert result["layout"] == "text_image"
        assert result["image_position"] == "right"

    def test_text_image_left_position(self, pptx_file, sample_image):
        result = json.loads(add_slide(
            path=pptx_file, title="Left Image", layout="text_image",
            text="Description", image=sample_image, image_position="left",
        ))
        assert result["image_position"] == "left"

    def test_bullets_image_layout(self, pptx_file, sample_image):
        result = json.loads(add_slide(
            path=pptx_file, title="Bullets+Image", layout="bullets_image",
            bullets=["A", "B"], image=sample_image,
        ))
        assert result["layout"] == "bullets_image"
        assert result["bullet_count"] == 2

    def test_two_column_layout(self, pptx_file):
        result = json.loads(add_slide(
            path=pptx_file, title="Two Column", layout="two_column",
            left_column=["L1", "L2"], right_column=["R1", "R2"],
        ))
        assert result["layout"] == "two_column"

    def test_two_column_json_strings(self, pptx_file):
        result = json.loads(add_slide(
            path=pptx_file, title="Two Col JSON", layout="two_column",
            left_column='["L1"]', right_column='["R1"]',
        ))
        assert result["layout"] == "two_column"

    def test_blank_layout(self, pptx_file):
        result = json.loads(add_slide(path=pptx_file, layout="blank"))
        assert result["layout"] == "blank"

    def test_unknown_layout(self, pptx_file):
        result = json.loads(add_slide(
            path=pptx_file, title="Bad", layout="unknown",
        ))
        assert "error" in result
        assert "valid_layouts" in result

    def test_nonexistent_file(self, tmp_data_dir):
        result = json.loads(add_slide(
            path=str(tmp_data_dir / "nope.pptx"), title="Fail", layout="text",
        ))
        assert "error" in result


class TestAddTableSlide:
    """Tests for add_table_slide()."""

    def test_basic_table(self, pptx_file):
        data = [["Name", "Age"], ["Alice", "30"], ["Bob", "25"]]
        result = json.loads(add_table_slide(
            path=pptx_file, title="Data Table", data=data,
        ))
        assert "table_size" in result
        assert result["table_size"] == "3 rows x 2 columns"

    def test_table_json_string(self, pptx_file):
        data_str = '[["H1", "H2"], ["V1", "V2"]]'
        result = json.loads(add_table_slide(
            path=pptx_file, title="JSON Table", data=data_str,
        ))
        assert result["table_size"] == "2 rows x 2 columns"

    def test_table_no_header_styling(self, pptx_file):
        data = [["A", "B"], ["1", "2"]]
        result = json.loads(add_table_slide(
            path=pptx_file, title="No Header", data=data, header=False,
        ))
        assert "error" not in result

    def test_empty_data(self, pptx_file):
        result = json.loads(add_table_slide(
            path=pptx_file, title="Empty", data=[],
        ))
        assert "error" in result

    def test_none_data(self, pptx_file):
        result = json.loads(add_table_slide(
            path=pptx_file, title="None Data",
        ))
        assert "error" in result


class TestAddImagesSlide:
    """Tests for add_images_slide()."""

    def test_horizontal_grid(self, pptx_file, sample_image):
        result = json.loads(add_images_slide(
            path=pptx_file, title="Images",
            images=[sample_image, sample_image], grid="horizontal",
        ))
        assert result["image_count"] == 2
        assert result["grid"] == "horizontal"

    def test_vertical_grid(self, pptx_file, sample_image):
        result = json.loads(add_images_slide(
            path=pptx_file, title="Vertical",
            images=[sample_image], grid="vertical",
        ))
        assert result["grid"] == "vertical"

    def test_2x2_grid(self, pptx_file, sample_image):
        result = json.loads(add_images_slide(
            path=pptx_file, title="Grid",
            images=[sample_image] * 4, grid="grid",
        ))
        assert result["image_count"] == 4

    def test_json_string_images(self, pptx_file, sample_image):
        images_json = json.dumps([sample_image, sample_image])
        result = json.loads(add_images_slide(
            path=pptx_file, title="JSON", images=images_json,
        ))
        assert result["image_count"] == 2

    def test_no_valid_images(self, pptx_file):
        result = json.loads(add_images_slide(
            path=pptx_file, title="Bad",
            images=["/nonexistent/a.png", "/nonexistent/b.png"],
        ))
        assert "error" in result

    def test_mixed_valid_invalid(self, pptx_file, sample_image):
        result = json.loads(add_images_slide(
            path=pptx_file, title="Mixed",
            images=[sample_image, "/nonexistent/bad.png"],
        ))
        assert result["image_count"] == 1


class TestStyleSlide:
    """Tests for style_slide()."""

    def test_named_color(self, pptx_file):
        result = json.loads(style_slide(
            path=pptx_file, background="blue",
        ))
        assert "error" not in result
        assert result["background"] == "blue"

    def test_hex_color(self, pptx_file):
        result = json.loads(style_slide(
            path=pptx_file, background="#FF5733",
        ))
        assert result["background"] == "#FF5733"

    def test_specific_slide_index(self, pptx_file):
        result = json.loads(style_slide(
            path=pptx_file, background="red", slide_index=0,
        ))
        assert result["slide_index"] == 0

    def test_negative_index(self, pptx_file):
        result = json.loads(style_slide(
            path=pptx_file, background="green", slide_index=-1,
        ))
        assert "error" not in result


class TestGetPresentationInfo:
    """Tests for get_presentation_info()."""

    def test_basic_info(self, pptx_file):
        result = json.loads(get_presentation_info(path=pptx_file))
        assert result["slide_count"] == 1
        assert result["width_inches"] == 10.0
        assert result["height_inches"] == 7.5
        assert len(result["slides"]) == 1
        assert result["slides"][0]["title"] == "Test Presentation"

    def test_nonexistent_file(self, tmp_data_dir):
        result = json.loads(get_presentation_info(
            path=str(tmp_data_dir / "nope.pptx"),
        ))
        assert "error" in result

    def test_after_adding_slides(self, pptx_file):
        add_slide(path=pptx_file, title="Slide 2", layout="text", text="Content")
        add_slide(path=pptx_file, title="Slide 3", layout="blank")
        result = json.loads(get_presentation_info(path=pptx_file))
        assert result["slide_count"] == 3


class TestReadPresentation:
    """Tests for read_presentation()."""

    def test_read_title_slide(self, pptx_file):
        result = json.loads(read_presentation(path=pptx_file))
        assert result["status"] == "success"
        assert result["slide_count"] == 1
        assert result["slides"][0]["title"] == "Test Presentation"

    def test_read_with_text_content(self, pptx_file):
        add_slide(path=pptx_file, title="Text", layout="text", text="Hello World")
        result = json.loads(read_presentation(path=pptx_file))
        slide = result["slides"][1]
        # At least one shape should contain "Hello World"
        texts = [s.get("text", "") for s in slide["shapes"]]
        assert any("Hello World" in t for t in texts)

    def test_read_with_table(self, pptx_file):
        add_table_slide(
            path=pptx_file, title="Table",
            data=[["H1", "H2"], ["V1", "V2"]],
        )
        result = json.loads(read_presentation(path=pptx_file))
        tables = result["slides"][1]["tables"]
        assert len(tables) == 1
        assert tables[0]["rows"][0] == ["H1", "H2"]

    def test_nonexistent_file(self, tmp_data_dir):
        result = json.loads(read_presentation(
            path=str(tmp_data_dir / "nope.pptx"),
        ))
        assert "error" in result


class TestModifyPresentation:
    """Tests for modify_presentation()."""

    def test_update_shape_text(self, pptx_file):
        result = json.loads(modify_presentation(
            path=pptx_file,
            updates=[{"slide_index": 0, "shape_index": 0, "text": "New Title"}],
        ))
        assert result["status"] == "success"
        assert result["shapes_updated"] == 1

        # Verify the change
        prs = Presentation(pptx_file)
        assert prs.slides[0].shapes[0].text_frame.text == "New Title"

    def test_update_json_string(self, pptx_file):
        updates_json = '[{"slide_index": 0, "shape_index": 0, "text": "JSON Updated"}]'
        result = json.loads(modify_presentation(
            path=pptx_file, updates=updates_json,
        ))
        assert result["shapes_updated"] == 1

    def test_delete_slide(self, pptx_file):
        add_slide(path=pptx_file, title="Extra", layout="text", text="To delete")
        prs = Presentation(pptx_file)
        assert len(prs.slides) == 2

        result = json.loads(modify_presentation(
            path=pptx_file, delete_slides=[1],
        ))
        assert result["slides_deleted"] == 1

    def test_no_updates_or_deletes(self, pptx_file):
        result = json.loads(modify_presentation(path=pptx_file))
        assert "error" in result

    def test_invalid_slide_index(self, pptx_file):
        result = json.loads(modify_presentation(
            path=pptx_file,
            updates=[{"slide_index": 99, "shape_index": 0, "text": "Fail"}],
        ))
        assert result["shapes_updated"] == 0

    def test_nonexistent_file(self, tmp_data_dir):
        result = json.loads(modify_presentation(
            path=str(tmp_data_dir / "nope.pptx"),
            updates=[{"slide_index": 0, "shape_index": 0, "text": "Fail"}],
        ))
        assert "error" in result
